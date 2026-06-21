"""Deterministic PPTX generator.

The model supplies a *structured outline* (title + slides, each with a heading and
bullet points). This module renders that outline into a real ``.pptx`` using a
fixed, code-controlled template — the model never writes rendering code.

Design goals:
- **Deterministic / dependency-light**: only ``python-pptx`` + the standard library.
- **Defensive**: tolerate loose model output (missing fields, wrong types, oversized
  payloads) by normalizing before rendering, so a malformed outline degrades
  gracefully instead of crashing the request.
- **Generic / open-source safe**: no company names, internal services, or private
  identifiers. The accent color and footer are overridable via environment variables
  with neutral defaults.

Public API:
    normalize_outline(raw) -> dict          # validate + clamp into a canonical shape
    build_pptx(outline) -> bytes            # render canonical outline → .pptx bytes
    safe_filename(title, ext="pptx") -> str # derive a safe download filename
"""

from __future__ import annotations

import io
import os
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Limits (clamp model output so a runaway outline can't blow up render) ----
MAX_SLIDES = 30
MAX_BULLETS_PER_SLIDE = 12
MAX_TITLE_CHARS = 120
MAX_BULLET_CHARS = 240
MAX_SUBTITLE_CHARS = 200
MAX_NOTES_CHARS = 2000

# --- Theme (env-overridable, neutral defaults — keep this module generic) ------
_DEFAULT_ACCENT = "008042"  # single accent; aligns with the product's primary green


def _accent_color() -> RGBColor:
    raw = (os.environ.get("PPTX_ACCENT_COLOR") or _DEFAULT_ACCENT).strip().lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", raw):
        raw = _DEFAULT_ACCENT
    return RGBColor.from_string(raw.upper())


_INK = RGBColor.from_string("1A1A1A")
_BODY = RGBColor.from_string("3C3C3C")
_MUTED = RGBColor.from_string("8A8A8A")
_WHITE = RGBColor.from_string("FFFFFF")

# 16:9 widescreen canvas
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN = Inches(0.9)


def _clean_text(value: Any, limit: int) -> str:
    """Coerce arbitrary model output to a single trimmed string, length-capped."""
    if value is None:
        return ""
    if not isinstance(value, str):
        value = str(value)
    # collapse internal whitespace runs; strip control chars that break XML
    value = "".join(ch for ch in value if ch == "\n" or ch >= " ")
    value = re.sub(r"[ \t]+", " ", value).strip()
    if len(value) > limit:
        value = value[: limit - 1].rstrip() + "…"
    return value


def _normalize_bullets(raw: Any) -> list[dict[str, Any]]:
    """Bullets may arrive as list[str] or list[{text, level}]; normalize to the
    latter with a clamped indent level (0 or 1 for MVP)."""
    if not isinstance(raw, list):
        if isinstance(raw, str) and raw.strip():
            raw = [raw]
        else:
            return []
    out: list[dict[str, Any]] = []
    for item in raw:
        if isinstance(item, dict):
            text = _clean_text(item.get("text") or item.get("content"), MAX_BULLET_CHARS)
            try:
                level = int(item.get("level") or 0)
            except (TypeError, ValueError):
                level = 0
        else:
            text = _clean_text(item, MAX_BULLET_CHARS)
            level = 0
        if not text:
            continue
        out.append({"text": text, "level": max(0, min(level, 1))})
        if len(out) >= MAX_BULLETS_PER_SLIDE:
            break
    return out


def normalize_outline(raw: Any) -> dict[str, Any]:
    """Validate + clamp loose model output into a canonical outline.

    Canonical shape::

        {
          "title": str,                 # always non-empty (falls back)
          "subtitle": str,              # may be ""
          "slides": [                   # 1..MAX_SLIDES
            {"title": str, "bullets": [{"text": str, "level": 0|1}], "notes": str},
            ...
          ],
        }

    Never raises on bad input — returns a renderable (possibly minimal) outline.
    """
    if not isinstance(raw, dict):
        raw = {}

    title = _clean_text(raw.get("title") or raw.get("topic"), MAX_TITLE_CHARS) or "未命名演示文稿"
    subtitle = _clean_text(raw.get("subtitle") or raw.get("author") or "", MAX_SUBTITLE_CHARS)

    raw_slides = raw.get("slides")
    if not isinstance(raw_slides, list):
        raw_slides = []

    slides: list[dict[str, Any]] = []
    for s in raw_slides:
        if not isinstance(s, dict):
            # tolerate a bare string as a section heading
            if isinstance(s, str) and s.strip():
                slides.append({"title": _clean_text(s, MAX_TITLE_CHARS), "bullets": [], "notes": ""})
            continue
        s_title = _clean_text(s.get("title") or s.get("heading"), MAX_TITLE_CHARS)
        bullets = _normalize_bullets(s.get("bullets") or s.get("points") or s.get("content"))
        notes = _clean_text(s.get("notes") or s.get("speaker_notes") or "", MAX_NOTES_CHARS)
        if not s_title and not bullets:
            continue
        if not s_title:
            s_title = "—"
        slides.append({"title": s_title, "bullets": bullets, "notes": notes})
        if len(slides) >= MAX_SLIDES:
            break

    if not slides:
        # guarantee at least one renderable content slide
        slides = [{"title": title, "bullets": [], "notes": ""}]

    return {"title": title, "subtitle": subtitle, "slides": slides}


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------


def _add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _add_rect(slide, left, top, width, height, color: RGBColor):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _render_title_slide(prs, accent: RGBColor, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # left accent band
    _add_rect(slide, Emu(0), Emu(0), Inches(0.35), _SLIDE_H, accent)

    _, tf = _add_textbox(slide, _MARGIN, Inches(2.4), _SLIDE_W - _MARGIN - Inches(0.6), Inches(2.4))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = _INK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(14)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        r2.font.color.rgb = _MUTED


def _render_content_slide(
    prs, accent: RGBColor, index: int, total: int, slide_data: dict[str, Any], footer: str
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # heading
    _, head_tf = _add_textbox(slide, _MARGIN, Inches(0.55), _SLIDE_W - _MARGIN - Inches(0.6), Inches(1.0))
    hp = head_tf.paragraphs[0]
    hr = hp.add_run()
    hr.text = slide_data["title"]
    hr.font.size = Pt(28)
    hr.font.bold = True
    hr.font.color.rgb = _INK

    # accent underline beneath the heading
    _add_rect(slide, _MARGIN, Inches(1.5), Inches(1.1), Pt(4), accent)

    # body bullets
    bullets = slide_data.get("bullets") or []
    if bullets:
        _, body_tf = _add_textbox(
            slide, _MARGIN, Inches(1.85), _SLIDE_W - _MARGIN - Inches(0.6), Inches(4.9)
        )
        body_tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, b in enumerate(bullets):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.level = b["level"]
            p.space_after = Pt(8)
            # manual bullet glyph (blank layout has no list formatting)
            marker = "•  " if b["level"] == 0 else "–  "
            mrun = p.add_run()
            mrun.text = marker
            mrun.font.size = Pt(18 if b["level"] == 0 else 16)
            mrun.font.color.rgb = accent
            mrun.font.bold = True
            trun = p.add_run()
            trun.text = b["text"]
            trun.font.size = Pt(18 if b["level"] == 0 else 16)
            trun.font.color.rgb = _BODY
    else:
        _, body_tf = _add_textbox(
            slide, _MARGIN, Inches(1.85), _SLIDE_W - _MARGIN - Inches(0.6), Inches(4.9)
        )
        p = body_tf.paragraphs[0]
        r = p.add_run()
        r.text = ""

    # footer: page number (+ optional footer text)
    _, foot_tf = _add_textbox(slide, _MARGIN, _SLIDE_H - Inches(0.55), _SLIDE_W - 2 * _MARGIN, Inches(0.35))
    fp = foot_tf.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"{footer} · {index}/{total}" if footer else f"{index}/{total}"
    fr.font.size = Pt(10)
    fr.font.color.rgb = _MUTED

    # speaker notes
    notes = slide_data.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_pptx(outline: dict[str, Any]) -> bytes:
    """Render a (canonical or loose) outline into ``.pptx`` bytes.

    Always normalizes first, so callers may pass raw model output directly.
    """
    outline = normalize_outline(outline)
    accent = _accent_color()
    footer = _clean_text(os.environ.get("PPTX_FOOTER_TEXT") or "", 60)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _render_title_slide(prs, accent, outline["title"], outline.get("subtitle", ""))

    slides = outline["slides"]
    total = len(slides)
    for i, s in enumerate(slides, start=1):
        _render_content_slide(prs, accent, i, total, s, footer)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_SLUG_RE = re.compile(r"[^\w一-鿿\- ]+", re.UNICODE)


def safe_filename(title: str, ext: str = "pptx") -> str:
    """Derive a filesystem/Content-Disposition-safe filename from a title.

    Keeps CJK + word chars, collapses the rest, caps length, guarantees an ext.
    """
    base = _SLUG_RE.sub("", title or "").strip()
    base = re.sub(r"\s+", " ", base)
    if not base:
        base = "presentation"
    if len(base) > 60:
        base = base[:60].rstrip()
    return f"{base}.{ext.lstrip('.')}"


if __name__ == "__main__":  # pragma: no cover — local smoke test
    import sys

    sample = {
        "title": "2024 Q4 营收分析",
        "subtitle": "示例副标题 · 2026-06",
        "slides": [
            {
                "title": "整体表现",
                "bullets": [
                    "Q4 总营收同比 +18.4%,创历史新高",
                    {"text": "其中线上渠道贡献 32%", "level": 1},
                    "门店数净增 47 家",
                ],
                "notes": "强调同比口径,排除新店爬坡期影响。",
            },
            {
                "title": "区域拆解",
                "bullets": ["华东领跑,占比 41%", "华南环比增速最高 +12%", "华北单店均值待提升"],
            },
            {"title": "下一步", "bullets": ["加密华南布点", "华北门店运营诊断", "线上会员复购专项"]},
        ],
    }
    data = build_pptx(sample)
    out = sys.argv[1] if len(sys.argv) > 1 else "/tmp/sample_outline.pptx"
    with open(out, "wb") as f:
        f.write(data)
    # re-open to validate the file is well-formed
    reopened = Presentation(io.BytesIO(data))
    print(f"OK: {len(data)} bytes, {len(reopened.slides.__iter__.__self__._sldIdLst)} slides → {out}")
    print(f"filename: {safe_filename(sample['title'])}")
